"""
exporters.py - Export data to Excel, PDF, PPTX.
Always receives connection, filters, vista, and output path.
Returns dict with result info.
"""

import sqlite3
from pathlib import Path

import queries


def excel(conn: sqlite3.Connection, filtros: dict, ruta: str, vista: str | None = None) -> dict:
    """Export filtered data to Excel."""
    import pandas as pd

    if not vista:
        exp = queries.explorar(conn)
        vista = exp.get("vista_principal", queries.VISTA_PREDETERMINADA)

    df = queries._cargar_dataframe(conn, vista)
    df = queries._aplicar_filtros(df, filtros)

    ruta_path = Path(ruta)
    ruta_path.parent.mkdir(parents=True, exist_ok=True)

    with pd.ExcelWriter(ruta, engine="openpyxl") as writer:
        df.to_excel(writer, sheet_name="Reporte", index=False)

        worksheet = writer.sheets["Reporte"]
        for column in worksheet.columns:
            max_length = max(len(str(cell.value or "")) for cell in column)
            worksheet.column_dimensions[column[0].column_letter].width = min(max_length + 2, 50)

    return {
        "archivo": str(ruta_path.resolve()),
        "registros": len(df),
    }


def pdf(conn: sqlite3.Connection, filtros: dict, ruta: str, vista: str | None = None) -> dict:
    """Export filtered data to PDF table using fpdf2."""
    from fpdf import FPDF

    import pandas as pd

    if not vista:
        exp = queries.explorar(conn)
        vista = exp.get("vista_principal", queries.VISTA_PREDETERMINADA)

    df = queries._cargar_dataframe(conn, vista)
    df = queries._aplicar_filtros(df, filtros)

    ruta_path = Path(ruta)
    ruta_path.parent.mkdir(parents=True, exist_ok=True)

    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Reporte de Contrataciones", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("Helvetica", "", 8)
    pdf.cell(0, 6, f"Registros: {len(df)}", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(4)

    cols = list(df.columns)
    col_w = max(15, int(270 / max(len(cols), 1)))

    pdf.set_font("Helvetica", "B", 7)
    for col in cols:
        pdf.cell(col_w, 6, str(col)[:col_w // 2], border=1, align="C")
    pdf.ln()

    pdf.set_font("Helvetica", "", 6)
    for _, row in df.head(200).iterrows():
        for col in cols:
            val = str(row.get(col, ""))[:col_w // 2]
            pdf.cell(col_w, 5, val, border=1, align="C")
        pdf.ln()
        if pdf.get_y() > 185:
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 7)
            for col in cols:
                pdf.cell(col_w, 6, str(col)[:col_w // 2], border=1, align="C")
            pdf.ln()
            pdf.set_font("Helvetica", "", 6)

    pdf.output(ruta)
    return {
        "archivo": str(ruta_path.resolve()),
        "registros": len(df),
    }


def pptx(conn: sqlite3.Connection, filtros: dict, ruta: str, vista: str | None = None) -> dict:
    """Export filtered data to PPTX table."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    import pandas as pd

    if not vista:
        exp = queries.explorar(conn)
        vista = exp.get("vista_principal", queries.VISTA_PREDETERMINADA)

    df = queries._cargar_dataframe(conn, vista)
    df = queries._aplicar_filtros(df, filtros)

    ruta_path = Path(ruta)
    ruta_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title text box (blank layout has no title placeholder)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    tf.text = "Reporte de Contrataciones"

    rows = min(len(df) + 1, 201)
    cols = len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5))
    table = table_shape.table

    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True

    for i, (_, row) in enumerate(df.head(200).iterrows()):
        for j, col in enumerate(df.columns):
            cell = table.cell(i + 1, j)
            cell.text = str(row.get(col, ""))
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)

    prs.save(ruta)
    return {
        "archivo": str(ruta_path.resolve()),
        "registros": len(df),
    }
